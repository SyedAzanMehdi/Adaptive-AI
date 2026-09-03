"""Build the professional, animated pitch deck: docs/17_Pro_Pitch_Deck.pptx.

18 slides, monochrome house style, built to impress a technical team lead:
product depth + engineering rigor + shipped evidence.

Real PowerPoint motion is injected as raw OOXML (python-pptx has no animation API):
  * a fade <p:transition> on every slide, and
  * a fade <p:timing> entrance sequence (build-on-click) on each slide's hero shapes.

Run from repo root:
    python scripts/build_pro_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# --- palette (monochrome) ---------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
INK = RGBColor(0x14, 0x14, 0x14)
GRAY_D = RGBColor(0x3D, 0x3D, 0x3D)
GRAY_M = RGBColor(0x6A, 0x6A, 0x6A)
GRAY_S = RGBColor(0x9A, 0x9A, 0x9A)
DIM = RGBColor(0xC6, 0xC6, 0xC6)      # light gray text on dark
LINE = RGBColor(0xDC, 0xDC, 0xDC)
PANEL = RGBColor(0xF4, 0xF4, 0xF4)
WATERMARK = RGBColor(0x24, 0x24, 0x24)
FONT = "Calibri"

SW, SH = 13.333, 7.5
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

ATTR = "Created by Syed Azan Mehdi Shah"
SUP = "Supervisor: Mr. Asif Raza — Lecturer, University of Mianwali"
BRAND = "Adaptive AI Learning Platform · University of Mianwali"

TOTAL = 18


# --- primitives -------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def add_slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK if dark else WHITE
    return s


def rect(slide, x, y, w, h, fill=BLACK, line_color=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    return sh


def para(tf, text, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space_before=0, space_after=6, first=False,
         line_spacing=1.0):
    if hasattr(tf, "text_frame"):
        tf = tf.text_frame
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = FONT
    return p


def card(slide, x, y, w, h, head, body, dark=False, head_size=15, body_size=11.5):
    sh = rect(slide, x, y, w, h,
              fill=(RGBColor(0x16, 0x16, 0x16) if dark else WHITE),
              line_color=(RGBColor(0x33, 0x33, 0x33) if dark else LINE), line_w=1.0)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.20)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.14)
    para(tf, head, size=head_size, color=(WHITE if dark else INK), bold=True,
         first=True, space_after=5)
    if body:
        para(tf, body, size=body_size, color=(DIM if dark else GRAY_D),
             space_after=0, line_spacing=1.06)
    return sh


def stat(slide, x, y, w, h, num, label, dark=False, num_size=34):
    sh = rect(slide, x, y, w, h,
              fill=(RGBColor(0x16, 0x16, 0x16) if dark else PANEL),
              line_color=(RGBColor(0x33, 0x33, 0x33) if dark else None), line_w=1.0)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    para(tf, num, size=num_size, color=(WHITE if dark else INK), bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=3)
    para(tf, label, size=10.5, color=(DIM if dark else GRAY_M),
         align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.0)
    return sh


def footer(slide, dark=False, n=None):
    c = GRAY_S if dark else GRAY_M
    f = textbox(slide, 0.7, 7.03, 9.5, 0.34)
    para(f, BRAND, size=9.5, color=c, first=True, space_after=0)
    if n:
        g = textbox(slide, 11.0, 7.03, 1.63, 0.34)
        para(g, f"{n:02d} / {TOTAL}", size=9.5, color=c, align=PP_ALIGN.RIGHT,
             first=True, space_after=0)


def header(slide, kick, title_text, dark=False, n=None):
    """Kicker + title + accent bar + footer. Returns the animated hero shapes."""
    k = textbox(slide, 0.7, 0.5, 11.93, 0.34)
    para(k, kick.upper(), size=12, color=(GRAY_S if dark else GRAY_M), bold=True,
         first=True, space_after=0)
    t = textbox(slide, 0.7, 0.85, 11.93, 0.92)
    para(t, title_text, size=29, color=(WHITE if dark else INK), bold=True,
         first=True, space_after=0)
    rect(slide, 0.72, 1.70, 0.95, 0.045, fill=(WHITE if dark else BLACK))
    footer(slide, dark=dark, n=n)
    return [k, t]


def bullets(slide, items, x=0.7, y=2.0, w=11.93, h=4.6, dark=False,
            size=15, gap=11):
    tb = textbox(slide, x, y, w, h)
    tf = tb.text_frame
    hc = WHITE if dark else INK
    bc = DIM if dark else GRAY_D
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            head, body = it
            para(tf, head, size=size, color=hc, bold=True, first=(i == 0), space_after=3)
            para(tf, body, size=size - 2.5, color=bc, space_after=gap, line_spacing=1.06)
        else:
            para(tf, "•   " + it, size=size, color=hc, first=(i == 0), space_after=gap)
    return tb


def col_header(slide, x, y, w, text):
    tb = textbox(slide, x, y, w, 0.34)
    para(tb, text.upper(), size=12, color=GRAY_M, bold=True, first=True, space_after=0)
    rect(slide, x + 0.01, y + 0.36, w - 0.02, 0.02, fill=LINE)
    return tb


# --- animation injection (raw OOXML) ---------------------------------------
def set_transition(slide):
    xml = (f'<p:transition xmlns:p="{P_NS}" spd="med" advClick="1">'
           f'<p:fade/></p:transition>')
    slide._element.append(etree.fromstring(xml))


def set_entrance(slide, shapes, effect_dur=450):
    """Fade-in entrance (build on click) for each shape, in order."""
    ids = [str(sh.shape_id) for sh in shapes]
    if not ids:
        return
    cid = [2]  # root=1, seq=2; effects start at 3

    def nid():
        cid[0] += 1
        return cid[0]

    inner = []
    for spid in ids:
        o, m, e, s1, s2 = nid(), nid(), nid(), nid(), nid()
        inner.append(f'''<p:par><p:cTn id="{o}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{m}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{e}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{s1}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set><p:animEffect transition="in" filter="fade"><p:cBhvr><p:cTn id="{s2}" dur="{effect_dur}"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>''')
    xml = (f'<p:timing xmlns:p="{P_NS}"><p:tnLst><p:par>'
           f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
           f'<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
           f'<p:childTnLst>{"".join(inner)}</p:childTnLst></p:cTn>'
           f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
           f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
           f'</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>')
    slide._element.append(etree.fromstring(xml))


def finish(slide, targets):
    set_transition(slide)
    set_entrance(slide, targets)


def arrow(slide, x, y, dark=False):
    tb = textbox(slide, x, y, 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tb, "→", size=22, color=(GRAY_S if dark else GRAY_S), bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    return tb


# --- slides -----------------------------------------------------------------
def build():
    prs = new_deck()

    # 1 — Title -------------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    k = textbox(s, 0.9, 1.15, 11.5, 0.4)
    para(k, "UNIVERSITY OF MIANWALI · PROJECT PITCH", size=13, color=GRAY_S,
         bold=True, first=True, space_after=0)
    tg.append(k)
    t = textbox(s, 0.9, 1.7, 11.5, 1.9)
    para(t, "Adaptive AI Learning Platform", size=46, color=WHITE, bold=True,
         first=True, space_after=12)
    para(t, "A live AI tutor for computer science that diagnoses each learner, rewrites "
            "lessons to their level, predicts what they will forget, and turns skill into "
            "scholarships, jobs, and first clients.", size=16, color=DIM, line_spacing=1.12)
    tg.append(t)
    rect(s, 0.95, 4.35, 1.25, 0.05, fill=WHITE)
    cr = textbox(s, 0.9, 4.6, 11.5, 1.0)
    para(cr, ATTR, size=13, color=RGBColor(0xB4, 0xB4, 0xB4), bold=True, first=True, space_after=4)
    para(cr, SUP, size=13, color=RGBColor(0xB4, 0xB4, 0xB4), space_after=0)
    tg.append(cr)
    band_y = 5.95
    for i, (num, lab) in enumerate([("14+", "student features shipped"),
                                    ("40/40", "automated tests green"),
                                    ("2 × 3", "AI key × model fallback")]):
        stat(s, 0.9 + i * 3.95, band_y, 3.7, 0.95, num, lab, dark=True, num_size=26)
    finish(s, tg)

    # 2 — Agenda ------------------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "Agenda", "What this deck covers", n=2)
    rows = [
        ("01", "The Problem", "Why static online CS education fails the people who need it most"),
        ("02", "The Solution", "A closed adaptive loop: diagnose, adapt, mentor, remember, convert"),
        ("03", "Signature Innovations", "Memory Twin, Struggle DNA, Career Autopilot, Design Dojo"),
        ("04", "Opportunity Layer", "Passport, Scholarships, Freelance, Compass, PathFinder"),
        ("05", "Technology & Reliability", "MERN + TypeScript monorepo, structured AI, security"),
        ("06", "Feasibility & Traction", "What is actually built, verified, and deployment-ready"),
    ]
    y = 2.05
    for num, head, sub in rows:
        nb = textbox(s, 0.7, y, 0.9, 0.6)
        para(nb, num, size=22, color=GRAY_S, bold=True, first=True, space_after=0)
        tb = textbox(s, 1.65, y + 0.02, 10.9, 0.72)
        para(tb, head, size=16, color=INK, bold=True, first=True, space_after=2)
        para(tb, sub, size=11.5, color=GRAY_M, space_after=0)
        rect(s, 0.7, y + 0.74, 11.93, 0.012, fill=LINE)
        tg += [nb, tb]
        y += 0.80
    finish(s, tg)

    # 3 — Divider 01 --------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    wm = textbox(s, 0.75, 0.7, 6.0, 3.2)
    para(wm, "01", size=170, color=WATERMARK, bold=True, first=True, space_after=0)
    tg.append(wm)
    tt = textbox(s, 0.95, 3.9, 11.0, 1.6)
    rect(s, 0.98, 3.75, 1.25, 0.05, fill=WHITE)
    para(tt, "The Problem", size=40, color=WHITE, bold=True, first=True, space_after=8)
    para(tt, "Static courses fail the people who need them most.", size=16, color=DIM)
    tg.append(tt)
    footer(s, dark=True, n=3)
    finish(s, tg)

    # 4 — Problem detail ----------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "01 · The Problem", "Three failures baked into static learning", n=4)
    cw = 3.83
    xs = [0.7, 0.7 + cw + 0.22, 0.7 + 2 * (cw + 0.22)]
    data = [
        ("One-size-fits-all", "Every learner gets identical content regardless of level or style, so most disengage and quietly drop out."),
        ("No memory of forgetting", "Platforms log what you completed, not what you retained. Skills decay silently and nobody intervenes in time."),
        ("No bridge to outcomes", "Learning ends at 'course complete' — no interview readiness, no proof for studying abroad, no first client."),
    ]
    for x, (h, b) in zip(xs, data):
        tg.append(card(s, x, 2.0, cw, 2.35, h, b))
    band = [
        ("≈90%", "of online-course learners never finish (industry estimates)"),
        ("1 path", "the same content for every level, style, and goal"),
        ("Next 1B", "emerging-market learners, underserved by US-priced products"),
    ]
    for i, (num, lab) in enumerate(band):
        stat(s, xs[i], 4.65, cw, 1.55, num, lab, num_size=30)
    finish(s, tg)

    # 5 — Who it affects ----------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "01 · Who It Affects", "Built for the learners incumbents overlook", n=5)
    cw, ch = 2.87, 3.15
    xs = [0.7 + i * (cw + 0.15) for i in range(4)]
    data = [
        ("Self-taught developers", "No curriculum, no feedback, no proof. They need a path and a mentor — not another video library."),
        ("University CS students", "Theory-heavy courses, little interview prep. They need applied skill and verifiable evidence."),
        ("Career switchers", "Coming from another field. They need a level-aware ramp and a portfolio, fast."),
        ("Emerging-market learners", "Price-sensitive, often non-native English. They need fair pricing and a local language."),
    ]
    for x, (h, b) in zip(xs, data):
        tg.append(card(s, x, 2.0, cw, ch, h, b, head_size=14, body_size=11))
    bl = textbox(s, 0.7, 5.5, 11.93, 0.9)
    para(bl, "Primary market: Pakistan, India, MENA, and Southeast Asia.  Secondary: universities "
             "and bootcamps (B2B2C) that need retention and outcomes evidence.",
         size=13, color=GRAY_D, bold=True, first=True, space_after=0, line_spacing=1.1)
    tg.append(bl)
    finish(s, tg)

    # 6 — Divider 02 --------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    wm = textbox(s, 0.75, 0.7, 6.0, 3.2)
    para(wm, "02", size=170, color=WATERMARK, bold=True, first=True, space_after=0)
    tg.append(wm)
    tt = textbox(s, 0.95, 3.9, 11.0, 1.6)
    rect(s, 0.98, 3.75, 1.25, 0.05, fill=WHITE)
    para(tt, "The Solution", size=40, color=WHITE, bold=True, first=True, space_after=8)
    para(tt, "A live AI tutor that adapts, remembers, and opens doors.", size=16, color=DIM)
    tg.append(tt)
    footer(s, dark=True, n=6)
    finish(s, tg)

    # 7 — The adaptive loop -------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "02 · The Solution", "One closed loop from first answer to first income", n=7)
    steps = [
        ("1 · Diagnose", "Adaptive, Gemini-generated tests build a live 5-domain Capability Matrix."),
        ("2 · Adapt", "Lessons rewrite to the learner's tier — analogies, diagrams, then internals."),
        ("3 · Mentor", "Code is critiqued on 4 axes: correctness, style, edge cases, optimization."),
        ("4 · Remember", "Memory Twin forecasts forgetting; Rescue Reviews reinforce before decay."),
        ("5 · Convert", "Skill becomes a Passport, scholarships, a freelance profile, a 90-day plan."),
    ]
    cw = 2.28
    gap = 0.135
    x = 0.7
    for i, (h, b) in enumerate(steps):
        tg.append(card(s, x, 2.15, cw, 2.6, h, b, head_size=13.5, body_size=10.5))
        if i < 4:
            arrow(s, x + cw - 0.02, 3.15)
        x += cw + gap
    bl = textbox(s, 0.7, 5.15, 11.93, 1.0)
    para(bl, "Every step feeds the same learner model — so the platform gets sharper, and more "
             "personal, the longer a student uses it.", size=13.5, color=GRAY_D, bold=True,
         first=True, space_after=0, line_spacing=1.1)
    tg.append(bl)
    finish(s, tg)

    # 8 — Adaptive engine ---------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "02 · The Engine", "How the adaptation actually works", n=8)
    tg.append(bullets(s, [
        ("Capability Matrix", "A five-domain proficiency model updated by every answer and submission — it drives all downstream adaptation."),
        ("Self-rewriting lessons", "Same learning objectives, personalized delivery; cached per concept × tier × style so repeats are instant."),
        ("4-axis code mentorship", "Tiered, structured feedback on correctness, style, edge cases, and optimization — student code is never executed on the server."),
        ("Domain-general AI mentor", "Persistent per-user chat with topic memory, graceful knowledge-base fallback, and strict history isolation."),
    ], y=2.0, size=15.5, gap=13))
    finish(s, tg)

    # 9 — Signature innovations --------------------------------------------
    s = add_slide(prs)
    tg = header(s, "03 · Signature Innovations", "Capabilities no incumbent ships", n=9)
    cw, ch = 5.9, 2.2
    positions = [(0.7, 2.0), (6.73, 2.0), (0.7, 4.35), (6.73, 4.35)]
    data = [
        ("Memory Twin™ — skill-decay prediction", "Fits forgetting curves R(t)=e^(−t/S) to real practice history; forecasts 14-day retention; 2-minute Rescue Reviews reinforce a skill before it fades."),
        ("Struggle DNA™ — cognitive phenotyping", "Mines behavior into Resilience, Depth Tolerance, Edge Awareness, and Craft; assigns a struggle archetype and prescribes targeted countermeasures."),
        ("Career Autopilot™ — JD to 90-day plan", "Paste any job description: importance-weighted gap analysis against the live matrix, a hire-readiness score, and a deterministic 90-day plan."),
        ("System Design Dojo™ — interview training", "Six structured challenges graded on a 4-axis interview rubric (Clarify → Estimate → Model → Architect → Scale) with AI critique and history."),
    ]
    for (x, y), (h, b) in zip(positions, data):
        tg.append(card(s, x, y, cw, ch, h, b, head_size=14, body_size=11))
    finish(s, tg)

    # 10 — Opportunity layer ------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "04 · Opportunity Layer", "Turning learning into real-world outcomes", n=10)
    cw, ch = 3.83, 2.5
    xs = [0.7, 0.7 + cw + 0.22, 0.7 + 2 * (cw + 0.22)]
    data = [
        ("Skill Passport™", "Portable, verifiable proof of skill: a stable passport ID, per-domain mastery, evidence counts, and issuer attestation. JSON download plus resume export."),
        ("Scholarship Radar™", "15 curated fully-funded international scholarships with live rolling-deadline countdowns, match scoring, and level/country/field filters. Free for all."),
        ("Freelance Launchpad™", "An AI-drafted, matrix-grounded freelance profile — niche, skills, starter gigs, and a realistic rate. It never invents skills the learner lacks."),
    ]
    for x, (h, b) in zip(xs, data):
        tg.append(card(s, x, 2.0, cw, ch, h, b, head_size=14, body_size=11))
    band = card(s, 0.7, 4.75, 11.93, 1.45, "Global Access Doctrine™ + planning tools",
                "PPP regional pricing (Pakistan ₨999 default across 8 markets) and an Urdu dual-language "
                "glossary make it reachable; Domain Compass™ maps 64 computing domains with 10-year demand "
                "trends, and PathFinder™ turns the live matrix into an adaptive 7-day study plan.",
                head_size=14, body_size=11.5)
    tg.append(band)
    finish(s, tg)

    # 11 — Divider 03 -------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    wm = textbox(s, 0.75, 0.7, 6.5, 3.2)
    para(wm, "03", size=170, color=WATERMARK, bold=True, first=True, space_after=0)
    tg.append(wm)
    tt = textbox(s, 0.95, 3.9, 11.0, 1.6)
    rect(s, 0.98, 3.75, 1.25, 0.05, fill=WHITE)
    para(tt, "Technology & Reliability", size=40, color=WHITE, bold=True, first=True, space_after=8)
    para(tt, "Engineering maturity a team lead can verify, not just admire.", size=16, color=DIM)
    tg.append(tt)
    footer(s, dark=True, n=11)
    finish(s, tg)

    # 12 — Architecture -----------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "05 · Architecture", "A clean, layered MERN + TypeScript monorepo", n=12)
    layers = [
        ("Client — React 19 + Vite", "Mobile-first monochrome SPA: the student experience and the admin console, with GSAP / Three.js / Motion and reduced-motion support."),
        ("API — Express 5, strict MVC", "Routes → controllers → services. Issuer-bound JWT, RBAC, plan gating, and Zod request validation at the edge."),
        ("AI services layer", "All Gemini orchestration is isolated here. Structured JSON outputs, a dual-key/model cascade, and a deterministic mock fallback."),
        ("Data — MongoDB / Mongoose", "Users, CapabilityMatrix, Lessons, CodeSubmission, ChatMessage, AuditLog — all validated against schemas shared with the client."),
    ]
    y = 2.0
    for h, b in layers:
        sh = rect(s, 0.7, y, 11.93, 1.02, fill=PANEL, line_color=LINE, line_w=1.0)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.24)
        tf.margin_right = Inches(0.22)
        para(tf, h, size=14.5, color=INK, bold=True, first=True, space_after=2)
        para(tf, b, size=11, color=GRAY_D, space_after=0, line_spacing=1.02)
        tg.append(sh)
        y += 1.12
    finish(s, tg)

    # 13 — AI reliability ---------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "05 · AI Reliability", "The demo can never fail on stage", n=13)
    tg.append(bullets(s, [
        ("Structured outputs everywhere", "Every AI call declares a response schema and is Zod-validated before anything is persisted — no free-form AI text in the database."),
        ("Dual-key, multi-model cascade", "On quota or overload, requests roll across 2 API keys × 3 models within a 45-second ceiling, honoring Google's retry hints."),
        ("Deterministic fallback", "No key, or total failure, drops to a mock provider that keeps every feature alive — graceful degradation instead of a 500."),
        ("Prompt-injection safe", "Student input is embedded strictly as data, never as instructions; all prompts are assembled server-side."),
    ], y=2.0, w=7.5, size=14, gap=11))
    stats = [("2 × 3", "key × model fallback cascade"), ("45s", "hard ceiling per AI request"), ("40/40", "automated tests green")]
    for i, (num, lab) in enumerate(stats):
        stat(s, 8.45, 2.0 + i * 1.42, 4.18, 1.25, num, lab, num_size=28)
    finish(s, tg)

    # 14 — Security & quality ----------------------------------------------
    s = add_slide(prs)
    tg = header(s, "05 · Security & Quality", "Production posture, verified", n=14)
    col_header(s, 0.7, 2.0, 5.9, "Security")
    tg.append(bullets(s, [
        "Issuer-bound JWT access + refresh tokens",
        "RBAC and plan gating (403 / 402) from verified claims",
        "Helmet CSP + HSTS; centralized security config",
        "Tiered rate limits on auth, chat, and every AI route",
        "Audit logging for admin mutations and plan grants",
        "NoSQL operator-injection guards; secrets only via env",
    ], x=0.7, y=2.5, w=5.85, h=4.1, size=13, gap=8))
    col_header(s, 6.78, 2.0, 5.85, "Quality")
    tg.append(bullets(s, [
        "40/40 Vitest + Supertest integration tests",
        "Full RBAC matrix and ownership coverage",
        "WCAG-AA contrast audit on every route, dark + light",
        "Responsive, mobile-first across breakpoints",
        "Persistent dev database; ephemeral, isolated tests",
        "Reproducible QA suite and signed-off QA report",
    ], x=6.78, y=2.5, w=5.85, h=4.1, size=13, gap=8))
    finish(s, tg)

    # 15 — Divider 04 -------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    wm = textbox(s, 0.75, 0.7, 6.5, 3.2)
    para(wm, "04", size=170, color=WATERMARK, bold=True, first=True, space_after=0)
    tg.append(wm)
    tt = textbox(s, 0.95, 3.9, 11.0, 1.6)
    rect(s, 0.98, 3.75, 1.25, 0.05, fill=WHITE)
    para(tt, "Feasibility & Traction", size=40, color=WHITE, bold=True, first=True, space_after=8)
    para(tt, "Not a concept deck — a shipped, running, tested product.", size=16, color=DIM)
    tg.append(tt)
    footer(s, dark=True, n=15)
    finish(s, tg)

    # 16 — What's actually built -------------------------------------------
    s = add_slide(prs)
    tg = header(s, "06 · What's Built", "Shipped end-to-end and running today", n=16)
    tiles = [("14+", "student features"), ("40/40", "tests passing"), ("5", "matrix domains"),
             ("64", "Compass domains"), ("15", "scholarships"), ("2", "AI keys cascaded")]
    tw = 1.87
    for i, (num, lab) in enumerate(tiles):
        stat(s, 0.7 + i * (tw + 0.14), 2.0, tw, 1.35, num, lab, num_size=26)
    col_header(s, 0.7, 3.7, 11.93, "Modules delivered")
    tg.append(bullets(s, [
        "Auth, RBAC, and premium plan gating — with a full admin console and audit log",
        "Adaptive diagnostics, self-rewriting lessons, code playground, and the AI mentor chat",
        "Memory Twin, Struggle DNA, Career Autopilot, and the System Design Dojo",
        "Skill Passport, Scholarship Radar, Freelance Launchpad, Domain Compass, and PathFinder",
        "Deployment-ready Vercel serverless adapter (same-origin, no CORS) awaiting production wiring",
    ], x=0.7, y=4.2, w=11.93, h=2.5, size=13.5, gap=9))
    finish(s, tg)

    # 17 — Roadmap & the ask ------------------------------------------------
    s = add_slide(prs)
    tg = header(s, "06 · Roadmap & Ask", "A clear, staged path from here", n=17)
    phases = [
        ("Now — Shipped", "Complete MVP running end-to-end, 40/40 tests, and a deployment-ready serverless build."),
        ("Next — Live", "Production deploy on Vercel + MongoDB Atlas, and Stripe billing for the premium tier."),
        ("Then — Retention", "Decay-alert notifications, streaks, and an institutional B2B2C dashboard for cohorts."),
        ("Later — Network", "Marketplace and university partnerships behind the Radar, Launchpad, and Passport."),
    ]
    cw, ch = 2.87, 2.9
    xs = [0.7 + i * (cw + 0.15) for i in range(4)]
    for x, (h, b) in zip(xs, phases):
        tg.append(card(s, x, 2.0, cw, ch, h, b, head_size=14, body_size=11))
    band = rect(s, 0.7, 5.15, 11.93, 1.05, fill=BLACK)
    tf = band.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.22)
    para(tf, "The ask", size=13, color=GRAY_S, bold=True, first=True, space_after=3)
    para(tf, "Hackathon validation, then a seed round to ship billing, scale the live deployment, "
             "and drive the emerging-market wedge.", size=13.5, color=WHITE, space_after=0, line_spacing=1.05)
    tg.append(band)
    finish(s, tg)

    # 18 — Closing ----------------------------------------------------------
    s = add_slide(prs, dark=True)
    tg = []
    t = textbox(s, 0.9, 2.0, 11.5, 2.4)
    para(t, "Thank you", size=52, color=WHITE, bold=True, first=True, space_after=12)
    para(t, "Adaptive AI Learning Platform — the tutor that never forgets you.",
         size=18, color=DIM, line_spacing=1.1)
    tg.append(t)
    rect(s, 0.95, 4.35, 1.25, 0.05, fill=WHITE)
    cr = textbox(s, 0.9, 4.6, 11.5, 1.4)
    para(cr, ATTR, size=13.5, color=RGBColor(0xB4, 0xB4, 0xB4), bold=True, first=True, space_after=4)
    para(cr, SUP, size=13.5, color=RGBColor(0xB4, 0xB4, 0xB4), space_after=10)
    para(cr, "Live demo and questions welcome.", size=13, color=GRAY_S, italic=True, space_after=0)
    tg.append(cr)
    footer(s, dark=True, n=18)
    finish(s, tg)

    prs.save("docs/17_Pro_Pitch_Deck.pptx")
    print(f"OK docs/17_Pro_Pitch_Deck.pptx ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
