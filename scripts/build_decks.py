"""Rebuild the pitch decks (05 Overview, 07 Investor, 09 Hackathon, 15 Competition Pitch, 16 University Pitch).

Monochrome black & white styling matching the product. Run from repo root:
    python scripts/build_decks.py            # build all decks
    python scripts/build_decks.py pitch      # build only docs/15
    python scripts/build_decks.py university # build only docs/16 (overview|investor|hackathon|pitch|university)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
GRAY_DARK = RGBColor(0x44, 0x44, 0x44)
GRAY_MID = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xE5, 0xE5, 0xE5)
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def line(tf, text, size=14, color=BLACK, bold=False, first=False, align=PP_ALIGN.LEFT, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def kicker(slide, text, dark=False):
    tf = box(slide, 0.7, 0.45, 12, 0.4)
    line(tf, text.upper(), size=12, color=WHITE if dark else GRAY_MID, bold=True, first=True)


def title(slide, text, dark=False):
    tf = box(slide, 0.7, 0.85, 12, 1.0)
    line(tf, text, size=30, color=WHITE if dark else BLACK, bold=True, first=True)
    bar = slide.shapes.add_shape(1, Inches(0.72), Inches(1.75), Inches(0.9), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE if dark else BLACK
    bar.line.fill.background()


def bullets(slide, items, x=0.7, y=2.0, w=11.9, h=5.0, size=15, dark=False, gap=10):
    tf = box(slide, x, y, w, h)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            head, body = item
            p = line(tf, head, size=size, color=WHITE if dark else BLACK, bold=True,
                     first=(i == 0), space_after=2)
            p2 = line(tf, body, size=size - 2, color=GRAY_LIGHT if dark else GRAY_DARK,
                      space_after=gap)
        else:
            line(tf, "•  " + item, size=size, color=WHITE if dark else BLACK,
                 first=(i == 0), space_after=gap)


def card(slide, x, y, w, h, head, body, inverted=False):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLACK if inverted else WHITE
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = head
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = WHITE if inverted else BLACK
    r.font.name = FONT
    p.space_after = Pt(4)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.size = Pt(11.5)
    r2.font.color.rgb = GRAY_LIGHT if inverted else GRAY_DARK
    r2.font.name = FONT


def footer(slide, text, dark=False):
    tf = box(slide, 0.7, 7.05, 12, 0.35)
    line(tf, text, size=10, color=GRAY_MID if not dark else RGBColor(0x99, 0x99, 0x99), first=True)


def title_slide(prs, kicker_text, main, sub, credit, credit2=None):
    s = blank(prs, BLACK)
    tf = box(s, 0.9, 1.5, 11.5, 0.5)
    line(tf, kicker_text.upper(), size=14, color=RGBColor(0xAA, 0xAA, 0xAA), bold=True, first=True)
    tf = box(s, 0.9, 2.1, 11.5, 2.2)
    line(tf, main, size=44, color=WHITE, bold=True, first=True, space_after=10)
    line(tf, sub, size=18, color=RGBColor(0xCC, 0xCC, 0xCC))
    bar = s.shapes.add_shape(1, Inches(0.95), Inches(4.9), Inches(1.2), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()
    tf = box(s, 0.9, 6.25 if credit2 else 6.6, 11.5, 0.9 if credit2 else 0.5)
    line(tf, credit, size=12, color=RGBColor(0x99, 0x99, 0x99), first=True)
    if credit2:
        line(tf, credit2, size=12, color=RGBColor(0x99, 0x99, 0x99))


ATTR = "Created by Syed Azan Mehdi Shah — AI-Driven Adaptive Learning Platform"


def build_overview():
    prs = new_deck()
    title_slide(
        prs,
        "AI-Driven Adaptive Education",
        "Adaptive AI Learning Platform",
        "Personalized computer-science education that predicts forgetting, profiles struggle, and turns learning into opportunity.",
        ATTR,
    )

    s = blank(prs)
    kicker(s, "The Core Problem")
    title(s, "Static courses fail almost everyone")
    bullets(s, [
        ("One-size-fits-all content", "Every learner sees the same lessons regardless of strength, weakness, or style — so most disengage."),
        ("No memory of forgetting", "Platforms record what you completed, not what you retained. Skills silently decay after assessment."),
        ("No path to outcomes", "Learning stops at 'course finished' — students still can't answer interviews, prove skills abroad, or land first clients."),
        ("Our answer", "A live AI tutor: adaptive diagnostics build a Capability Matrix, lessons rewrite themselves per learner, and an opportunity layer converts skill into scholarships, jobs, and income."),
    ])
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "Adaptive Engine")
    title(s, "The learning loop")
    bullets(s, [
        ("Adaptive diagnostics", "Gemini-generated questions adapt to every answer; correct answers never leave the server. Compiles a 5-domain Capability Matrix."),
        ("Self-rewriting lessons", "Analogies for beginners, diagrams for intermediates, internals for advanced — same objectives, personalized delivery, cached per concept × tier × style."),
        ("4-axis code mentorship", "Correctness, style, edge cases, optimization — tiered feedback that feeds the matrix and Struggle DNA with every submission."),
        ("Domain-general AI mentor", "Persistent per-user chat with topic tagging, graceful knowledge-base fallback, strict history isolation."),
    ])
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "Signature Innovations")
    title(s, "Capabilities no incumbent ships")
    card(s, 0.7, 2.0, 5.9, 2.3, "Memory Twin™ — skill-decay prediction",
         "Fits forgetting curves R(t)=e^(−t/S) to real practice history; 14-day retention forecast; 2-minute Rescue Reviews reinforce stability before skills fade.")
    card(s, 6.75, 2.0, 5.9, 2.3, "Struggle DNA™ — cognitive phenotyping",
         "Mines behavior into Resilience, Depth Tolerance, Edge Awareness, Craft; assigns a struggle archetype and prescribes targeted countermeasures.")
    card(s, 0.7, 4.5, 5.9, 2.3, "Career Autopilot™ — JD → 90-day plan",
         "Paste any job description: importance-weighted gap analysis against the live matrix, hire-readiness score, deterministic 90-day plan.")
    card(s, 6.75, 4.5, 5.9, 2.3, "System Design Dojo™ — interview training",
         "Six structured challenges graded on a 4-axis interview rubric (Clarify → Estimate → Model → Architect → Scale) with AI critique and history.")
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "New — Global Opportunity Layer")
    title(s, "From learning to outcomes")
    card(s, 0.7, 2.0, 3.85, 4.6, "Skill Passport™",
         "Portable, verifiable proof of skill: deterministic passport ID, per-domain mastery, evidence counts, issuer attestation. JSON download + resume export.")
    card(s, 4.74, 2.0, 3.85, 4.6, "Scholarship Radar™",
         "15 curated fully-funded international scholarships with live rolling deadline countdowns, match scoring, and level/country/field filters. Free for all.")
    card(s, 8.78, 2.0, 3.85, 4.6, "Freelance Launchpad™",
         "AI-drafted, matrix-grounded freelance profile — niche, skills, starter gigs, realistic rate. Copy straight to any marketplace; never invents skills.")
    tf = box(s, 0.7, 6.75, 12, 0.4)
    line(tf, "Plus the Global Access Doctrine™: PPP regional pricing (Pakistan ₨999 default, 8 markets) and an Urdu dual-language glossary.",
         size=12, color=GRAY_DARK, bold=True, first=True)
    footer(s, ATTR)

    s = blank(prs, BLACK)
    kicker(s, "Architecture & Reliability", dark=True)
    title(s, "Demo-proof by design", dark=True)
    bullets(s, [
        ("MERN + TypeScript monorepo", "React 19 + Vite, Express 5 strict MVC, Mongoose, shared Zod schemas; AI isolated in a dedicated services layer."),
        ("Structured outputs everywhere", "Every Gemini call declares a response schema and validates before persistence; deterministic mock fallback keeps every feature alive with zero API keys."),
        ("Production security", "Issuer-bound JWT, RBAC + plan gating (403/402), helmet CSP + HSTS, tiered rate limits, audit logging."),
        ("Verified quality", "40/40 automated tests, 52/52 live API QA checks, WCAG-AA contrast audit on every route in dark + light."),
    ], dark=True)
    footer(s, ATTR, dark=True)

    prs.save("docs/05_Platform_Overview.pptx")
    print("OK docs/05_Platform_Overview.pptx (6 slides)")


def build_investor():
    prs = new_deck()
    title_slide(
        prs,
        "Adaptive+ Investor Pitch",
        "The retention engine for online CS education",
        "Memory modeling + struggle phenotyping + an opportunity layer that converts learning into scholarships, jobs, and income.",
        ATTR,
    )

    s = blank(prs)
    kicker(s, "Market Opportunity")
    title(s, "A $300B+ gap with a retention crisis")
    bullets(s, [
        ("Massive churn", "Over 90% of students drop out of static video-based CS courses; completion is the industry's unsolved problem."),
        ("Emerging-market surge", "Pakistan, India, MENA, Southeast Asia hold the next billion learners — underserved by US-priced, English-only products."),
        ("Outcomes are the product", "Students pay for results: passing interviews, studying abroad, first freelance income. Nobody connects learning evidence to those outcomes."),
    ])
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "Product Differentiation")
    title(s, "Why Adaptive+ wins")
    card(s, 0.7, 2.0, 5.9, 2.3, "Memory Twin™ retention engine",
         "Students watch their own skills fade on a forecast chart — then subscribe to stop it. Loss-aversion monetization built into the product.")
    card(s, 6.75, 2.0, 5.9, 2.3, "Struggle DNA™ profiler",
         "Coaches WHY students keep failing, not just what they missed — the moat compounds with every answer recorded.")
    card(s, 0.7, 4.5, 5.9, 2.3, "Career Autopilot™ killer demo",
         "Paste any JD → hire-readiness score + 90-day plan. One-line demo, clear outcome metric, natural premium upsell.")
    card(s, 6.75, 4.5, 5.9, 2.3, "Demo-proof architecture",
         "Every AI path has a structured schema + deterministic fallback — the live demo can never fail on stage or in diligence.")
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "New — The Wedge Into Emerging Markets")
    title(s, "Global Access Doctrine + Opportunity Layer")
    bullets(s, [
        ("Fair pricing, native language", "PPP price points across 8 regional markets (₨999 Pakistan default) and an Urdu dual-language glossary — incumbents sell US-priced, English-only."),
        ("Scholarship Radar™", "15 fully-funded programmes with live deadline countdowns and match scoring — the study-abroad dream, systematized and free, driving acquisition."),
        ("Freelance Launchpad™", "Turns the Capability Matrix into a marketplace-ready profile with starter gigs and realistic rates — income is the strongest retention mechanic."),
        ("Skill Passport™", "Portable, verifiable evidence for university applications, visa skill proof, and remote roles — no incumbent issues it."),
    ])
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "Monetization & Traction Path")
    title(s, "Freemium that sells itself")
    bullets(s, [
        ("Explorer (Free)", "Diagnostic, adaptive lessons, playground, Dojo, Passport, Scholarships, Freelance — the opportunity layer is free by design: it acquires and retains."),
        ("Adaptive+ (Premium)", "Memory Twin™, Rescue Reviews, full Struggle DNA™, Career Autopilot™. Gated server-side via JWT plan claims (402), not UI tricks."),
        ("Conversion hook", "Free users watch their skills decay on a chart, then subscribe to stop it — an emotional event the product creates on its own."),
        ("Roadmap", "Stripe billing → Atlas deployment → decay-alert notifications → institutional B2B2C tier → marketplace partnerships behind the Radar and Launchpad."),
    ])
    footer(s, ATTR)

    s = blank(prs, BLACK)
    kicker(s, "Evidence", dark=True)
    title(s, "Shipped, tested, verified", dark=True)
    bullets(s, [
        ("Complete MVP in the repo", "14-feature student experience + full admin console, all running end-to-end today."),
        ("40/40 automated tests · 52/52 live API QA checks", "RBAC matrix, plan gating, AI fallbacks, ownership, hardening."),
        ("Security posture", "Issuer-bound JWT, RBAC + plan gating, helmet CSP/HSTS, tiered rate limits, audit logging."),
        ("Ask", "Hackathon validation → seed round to ship Stripe billing, real deployment, and the emerging-market wedge."),
    ], dark=True)
    footer(s, ATTR, dark=True)

    prs.save("docs/07_Investor_Pitch_Deck.pptx")
    print("OK docs/07_Investor_Pitch_Deck.pptx (6 slides)")


def build_hackathon():
    prs = new_deck()
    title_slide(
        prs,
        "Hackathon Presentation",
        "Adaptive AI — the tutor that never forgets you",
        "Predicts when you'll forget a skill, profiles how you struggle, and converts learning into scholarships, jobs, and first clients.",
        ATTR,
    )

    s = blank(prs)
    kicker(s, "Live Demo Flow")
    title(s, "Five minutes, six wow moments")
    bullets(s, [
        ("1 · Register → adaptive diagnostic", "Questions adapt live to every answer; watch the Capability Matrix radar build itself."),
        ("2 · Lessons rewrite themselves", "Drop a domain's mastery and the lesson re-renders for your level and style."),
        ("3 · Memory Twin™ + Struggle DNA™", "Show the decay forecast, run a 2-minute Rescue Review, reveal the struggle archetype."),
        ("4 · Career Autopilot™", "Paste a real JD → hire-readiness score + 90-day plan in seconds."),
        ("5 · Opportunity layer", "Skill Passport JSON export, Scholarship Radar countdowns, one-click Freelance profile."),
    ])
    footer(s, ATTR)

    s = blank(prs)
    kicker(s, "New Since the Last Review")
    title(s, "Global Opportunity Layer")
    card(s, 0.7, 2.0, 3.85, 4.3, "Scholarship Radar™",
         "15 fully-funded scholarships, live rolling deadline countdowns, deterministic match scoring, 3 filters. Built for Pakistan + Global South students.")
    card(s, 4.74, 2.0, 3.85, 4.3, "Freelance Launchpad™",
         "Reads the live matrix, drafts headline / niche / skills / gigs / rate via Gemini structured output; deterministic fallback never invents skills.")
    card(s, 8.78, 2.0, 3.85, 4.3, "Skill Passport™ + resume",
         "Stable AP-XXXX ID, mastery evidence, attestation; download JSON, copy to clipboard, or copy as a plain-text resume.")
    tf = box(s, 0.7, 6.5, 12, 0.6)
    line(tf, "All three are free for every learner, rate-limited, RBAC-checked, and covered by new automated tests.",
         size=12.5, color=GRAY_DARK, bold=True, first=True)
    footer(s, ATTR)

    s = blank(prs, BLACK)
    kicker(s, "Technical Standards", dark=True)
    title(s, "Built with rigor", dark=True)
    bullets(s, [
        ("MERN + TypeScript monorepo", "React 19 + Vite, Express 5 strict MVC, AI isolated in services/, shared Zod schemas."),
        ("Structured AI outputs + fallbacks", "Every Gemini call schema-constrained; deterministic mock keeps every feature demo-proof offline."),
        ("Security", "Issuer-bound JWT, RBAC + plan gating, helmet CSP/HSTS, tiered rate limits, prompt-injection-safe prompts."),
        ("Verification", "40/40 tests · 52/52 live API QA checks · WCAG-AA contrast audit on every route in both themes."),
    ], dark=True)
    footer(s, ATTR, dark=True)

    prs.save("docs/09_Hackathon_Pitch.pptx")
    print("OK docs/09_Hackathon_Pitch.pptx (4 slides)")


PITCH_SUB = ("A live AI tutor for computer science that diagnoses each learner, rewrites lessons "
             "to their level, predicts what they will forget, and turns skill into scholarships, "
             "jobs, and first clients.")


def _pitch_section_slides(prs, foot):
    """The five pitch sections: problem, solution, need/impact, innovation/tech, feasibility."""
    # 01 — The problem you are solving, and who it affects
    s = blank(prs)
    kicker(s, "01 · The Problem")
    title(s, "Static courses fail the people who need them most")
    bullets(s, [
        ("One-size-fits-all content", "Every learner sees the same lessons regardless of strength, weakness, or style — so most disengage. Over 90% of students drop out of static online CS courses."),
        ("No memory of forgetting", "Platforms record what you completed, not what you retained. Skills silently decay after the assessment, and nobody warns the learner in time."),
        ("Learning stops at 'course finished'", "Students still cannot pass interviews, prove skills for study abroad, or land a first client. There is no bridge from learning to real outcomes."),
        ("Who it affects", "Self-taught and university CS learners in price-sensitive, underserved markets — Pakistan, India, MENA, Southeast Asia — where US-priced, English-only products leave the next billion learners behind."),
    ])
    footer(s, foot)

    # 02 — Your solution, and the audience it serves
    s = blank(prs)
    kicker(s, "02 · The Solution")
    title(s, "A live AI tutor that adapts, remembers, and opens doors")
    card(s, 0.7, 2.0, 5.9, 2.15, "Diagnose → adapt → mentor",
         "Gemini-generated adaptive diagnostics build a 5-domain Capability Matrix; lessons rewrite themselves to each learner's level and style; code gets 4-axis tiered mentorship.")
    card(s, 6.75, 2.0, 5.9, 2.15, "It remembers you",
         "Memory Twin™ fits forgetting curves to real practice and forecasts 14-day retention; Struggle DNA™ profiles how a learner fails, not just what they missed.")
    card(s, 0.7, 4.3, 5.9, 2.15, "It plans the path",
         "Domain Compass™ maps 64 computing domains with 10-year demand trends; PathFinder™ turns the live matrix into an adaptive 7-day study plan.")
    card(s, 6.75, 4.3, 5.9, 2.15, "It converts skill into opportunity",
         "Skill Passport™ (verifiable proof), Scholarship Radar™ (15 fully-funded programmes, live deadlines), Freelance Launchpad™, and Career Autopilot™ (JD → 90-day plan).")
    tf = box(s, 0.7, 6.6, 12, 0.4)
    line(tf, "Audience: CS students and self-taught developers in emerging markets first; universities and bootcamps (B2B2C) second. The opportunity layer is free by design; the retention engine is the premium tier.",
         size=11.5, color=GRAY_DARK, bold=True, first=True)
    footer(s, foot)

    # 03 — The need it addresses and the impact it makes
    s = blank(prs)
    kicker(s, "03 · Need & Impact")
    title(s, "From completion to consequence")
    bullets(s, [
        ("The need", "Learners do not need more content — they need to know what to learn next, remember it, and prove it. Incumbents optimize watch-time; nobody owns retention or outcomes."),
        ("Impact on the learner", "A personalized path, an early warning before a skill fades, and portable evidence — turning months of scattered effort into interviews passed, scholarships won, and first income earned."),
        ("Impact on the market", "PPP regional pricing (Pakistan ₨999 default, 8 markets) and an Urdu dual-language glossary make quality CS education reachable where US-priced, English-only products never land."),
        ("Why it compounds", "Every answer feeds the Capability Matrix, Struggle DNA, and Memory Twin — the model of each learner sharpens with use, so retention and outcomes improve the longer they stay."),
    ])
    footer(s, foot)

    # 04 — The innovation and the technology behind it (two columns)
    s = blank(prs)
    kicker(s, "04 · Innovation & Technology")
    title(s, "Novel capabilities on a demo-proof stack")
    tf = box(s, 0.7, 1.95, 5.9, 0.35)
    line(tf, "THE INNOVATION", size=12, color=GRAY_MID, bold=True, first=True)
    bullets(s, [
        ("Memory Twin™", "Skill-decay prediction from forgetting curves + 2-minute Rescue Reviews before a skill fades."),
        ("Struggle DNA™", "Cognitive phenotyping — a struggle archetype with targeted countermeasures."),
        ("Career Autopilot™", "Paste any job description → importance-weighted gap analysis + a 90-day plan."),
        ("Global Access Doctrine™", "PPP pricing and an Urdu glossary — access, not just content."),
    ], x=0.7, y=2.35, w=5.9, h=4.4, size=13, gap=8)
    tf = box(s, 6.75, 1.95, 5.9, 0.35)
    line(tf, "THE TECHNOLOGY", size=12, color=GRAY_MID, bold=True, first=True)
    bullets(s, [
        ("MERN + TypeScript monorepo", "React 19 + Vite, Express 5 strict MVC, Mongoose, shared Zod schemas."),
        ("AI isolated + structured", "All Gemini calls live in a services layer, declare a response schema, validate before persisting."),
        ("Never fails live", "Dual-key/model cascade with quota + overload rollover and a deterministic mock fallback."),
        ("Secure by design", "Issuer-bound JWT, RBAC + plan gating, helmet CSP/HSTS, tiered rate limits, audit logging."),
    ], x=6.75, y=2.35, w=5.9, h=4.4, size=13, gap=8)
    footer(s, foot)

    # 05 — Feasibility, and what you have actually built
    s = blank(prs, BLACK)
    kicker(s, "05 · Feasibility & Traction", dark=True)
    title(s, "Not a slide — a shipped, running product", dark=True)
    bullets(s, [
        ("Built and running today", "A 14-feature student experience plus a full admin console run end-to-end locally on the monorepo above."),
        ("Verified quality", "40/40 automated tests (RBAC matrix, plan gating, AI fallbacks, ownership, hardening) and a WCAG-AA contrast audit on every route in dark + light."),
        ("Feasible to run and scale", "Free-tier infrastructure — MongoDB Atlas M0, Vercel serverless, Gemini free tier — with fallbacks that remove any single-API dependency; one maintainer can operate it."),
        ("Deployment-ready", "The Vercel single-project serverless adapter and config are complete (same-origin, no CORS); it goes live the moment Atlas and the deploy token are wired."),
        ("The ask", "Hackathon validation, then a seed round to ship Stripe billing, the live production deploy, and the emerging-market wedge."),
    ], dark=True)
    footer(s, foot, dark=True)


def build_pitch_deck():
    """Competition pitch: problem, solution, need/impact, innovation/tech, feasibility."""
    prs = new_deck()
    title_slide(prs, "Competition Pitch", "Adaptive AI Learning Platform", PITCH_SUB, ATTR)
    _pitch_section_slides(prs, ATTR)
    prs.save("docs/15_Pitch_Deck.pptx")
    print("OK docs/15_Pitch_Deck.pptx (6 slides)")


def build_university_deck():
    """University submission pitch — same five sections, credited with the supervisor."""
    supervisor = "Supervisor: Mr. Asif Raza — Lecturer, University of Mianwali"
    foot = "Syed Azan Mehdi Shah · Supervisor: Mr. Asif Raza, Lecturer — University of Mianwali"
    prs = new_deck()
    title_slide(
        prs,
        "University of Mianwali · Project Pitch",
        "Adaptive AI Learning Platform",
        PITCH_SUB,
        ATTR,
        supervisor,
    )
    _pitch_section_slides(prs, foot)
    prs.save("docs/16_University_Pitch_Deck.pptx")
    print("OK docs/16_University_Pitch_Deck.pptx (6 slides)")


if __name__ == "__main__":
    import sys

    builders = {
        "overview": build_overview,
        "investor": build_investor,
        "hackathon": build_hackathon,
        "pitch": build_pitch_deck,
        "university": build_university_deck,
    }
    which = sys.argv[1].lower() if len(sys.argv) > 1 else "all"
    if which == "all":
        for fn in builders.values():
            fn()
    elif which in builders:
        builders[which]()
    else:
        print(f"Unknown deck '{which}'. Choose from: {', '.join(builders)} (or 'all').")
        sys.exit(1)
